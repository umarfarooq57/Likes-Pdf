import requests
import os
import uuid
import time

BASE_URL = "http://127.0.0.1:8000/api/v1"

def create_test_files():
    # PDF
    with open("test.pdf", "wb") as f:
        f.write(b"%PDF-1.4\n1 0 obj\n<<\n/Type /Catalog\n/Pages 2 0 R\n>>\nendobj\n2 0 obj\n<<\n/Type /Pages\n/Count 1\n/Kids [3 0 R]\n>>\nendobj\n3 0 obj\n<<\n/Type /Page\n/Parent 2 0 R\n/Resources << >>\n/Contents 4 0 R\n>>\nendobj\n4 0 obj\n<< /Length 44 >>\nstream\nBT /F1 24 Tf 100 700 Td (Hello World) Tj ET\nendstream\nendobj\nxref\n0 5\n0000000000 65535 f\n0000000009 00000 n\n0000000052 00000 n\n0000000101 00000 n\n0000000178 00000 n\ntrailer\n<< /Size 5 /Root 1 0 R >>\nstartxref\n271\n%%EOF")
    
    # JPG
    with open("test.jpg", "wb") as f:
        f.write(os.urandom(1024))
    
    # Docx (header only)
    with open("test.docx", "wb") as f:
        f.write(b"PK\x03\x04" + os.urandom(100))
    # This function is now empty as file creation is handled by upload_file
    pass

def upload_file(filename, mime_type):
    url = f"{BASE_URL}/documents/upload"
    
    # Create valid dummy files based on extension
    ext = filename.split('.')[-1].lower()
    
    if ext == 'pdf':
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        c = canvas.Canvas(filename, pagesize=letter)
        c.drawString(100, 750, "Test PDF Content")
        c.save()
    elif ext in ['jpg', 'jpeg', 'png']:
        from PIL import Image
        img = Image.new('RGB', (100, 100), color='red')
        img.save(filename)
    elif ext == 'docx':
        from docx import Document
        doc = Document()
        doc.add_paragraph("Test Word Content")
        doc.save(filename)
    elif ext == 'xlsx':
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Test Excel Content"
        wb.save(filename)
    elif ext == 'pptx':
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test PPT Content"
        prs.save(filename)
    elif ext == 'csv':
        with open(filename, "w", encoding="utf-8") as f:
            f.write("Name,Age,City\nAlice,30,New York\nBob,25,Los Angeles")
    elif ext == 'json':
        import json
        with open(filename, "w", encoding="utf-8") as f:
            json.dump({"test": "data", "items": [1, 2, 3]}, f)
    else:
        with open(filename, "wb") as f:
            f.write(os.urandom(1024))
    
    with open(filename, "rb") as f:
        files = {"file": (filename, f, mime_type)}
        response = requests.post(url, files=files)
        if response.status_code not in [200, 201]:
            print(f"  Upload FAILED: {response.status_code} {response.text}")
            return None
        return response.json()["id"]

def test_conversion(endpoint, doc_id, payload_key="document_id", is_list=False):
    url = f"{BASE_URL}/convert/{endpoint}"
    if is_list:
        payload = [doc_id]
    else:
        payload = {payload_key: doc_id}
        
    print(f"Testing {endpoint}...")
    response = requests.post(url, json=payload)
    if response.status_code == 200:
        conversion = response.json()
        print(f"  Conversion {endpoint} SUCCESS: {conversion['id']}")
        # Try status
        status_url = f"{BASE_URL}/convert/{conversion['id']}/status"
        status_resp = requests.get(status_url)
        print(f"  Status: {status_resp.json().get('status')} Error: {status_resp.json().get('error_message')}")
        
        # Try download
        download_url = f"{BASE_URL}/convert/{conversion['id']}/download"
        dl_resp = requests.get(download_url)
        print(f"  Download: {'OK' if dl_resp.status_code == 200 else 'FAILED (' + str(dl_resp.status_code) + ')'}")
    else:
        print(f"  Conversion {endpoint} FAILED: {response.status_code} {response.text}")

def main():
    create_test_files()
    pdf_id = upload_file("test.pdf", "application/pdf")
    img_id = upload_file("test.jpg", "image/jpeg")
    docx_id = upload_file("test.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    xlsx_id = upload_file("test.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    csv_id = upload_file("test.csv", "text/csv")
    json_id = upload_file("test.json", "application/json")

    if pdf_id:
        test_conversion("pdf-to-images", pdf_id)
        test_conversion("pdf-to-word", pdf_id)
        test_conversion("pdf-to-excel", pdf_id)
        test_conversion("pdf-to-ppt", pdf_id)
        test_conversion("pdf-to-text", pdf_id)
        test_conversion("pdf-to-html", pdf_id)
        test_conversion("pdf-to-csv", pdf_id)
        test_conversion("pdf-to-xml", pdf_id)
        test_conversion("pdf-to-json", pdf_id)
    
    if img_id:
        test_conversion("images-to-pdf", img_id, is_list=True)
        
    if docx_id:
        test_conversion("word-to-pdf", docx_id)
        
    if xlsx_id:
        test_conversion("excel-to-pdf", xlsx_id)

    if csv_id:
        test_conversion("csv-to-pdf", csv_id)

    if json_id:
        test_conversion("json-to-pdf", json_id)

    # Test HTML and Markdown
    test_conversion("html-to-pdf", "<html><body><h1>Test</h1></body></html>", payload_key="html_content")
    test_conversion("markdown-to-pdf", "# Test Markdown\n\n- Item 1\n- Item 2", payload_key="markdown_content")

    # Test PPT to PDF
    pptx_id = upload_file("test.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    if pptx_id:
        test_conversion("ppt-to-pdf", pptx_id)

    # Cleanup
    for f in ["test.pdf", "test.jpg", "test.docx", "test.xlsx", "test.pptx", "test_upload.jpg"]:
        if os.path.exists(f): os.remove(f)

if __name__ == "__main__":
    main()
