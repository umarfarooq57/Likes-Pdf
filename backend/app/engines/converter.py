"""
Document Converter
Convert various document formats to PDF
Provides fallback stubs when dependencies are not available
"""

from pathlib import Path
from typing import Optional

# Optional imports with feature flags
WEASYPRINT_AVAILABLE = False
PYTHON_DOCX_AVAILABLE = False
OPENPYXL_AVAILABLE = False
PYTHON_PPTX_AVAILABLE = False
MARKDOWN_AVAILABLE = False
PILLOW_AVAILABLE = False

try:
    from weasyprint import HTML, CSS
    WEASYPRINT_AVAILABLE = True
except ImportError:
    pass

try:
    from docx import Document as WordDocument
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    pass

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    pass

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    pass


class DocumentConverter:
    """Convert various document formats to PDF"""
    
    @staticmethod
    def html_to_pdf(html_path: Path, output_path: Path, css_path: Optional[Path] = None) -> Path:
        """Convert HTML to PDF"""
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        stylesheets = [CSS(filename=str(css_path))] if css_path else None
        HTML(filename=str(html_path)).write_pdf(str(output_path), stylesheets=stylesheets)
        return output_path
    
    @staticmethod
    def html_string_to_pdf(html_string: str, output_path: Path) -> Path:
        """Convert HTML string to PDF"""
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        HTML(string=html_string).write_pdf(str(output_path))
        return output_path
    
    @staticmethod
    def markdown_to_pdf(md_path: Path, output_path: Path) -> Path:
        """Convert Markdown to PDF"""
        if not MARKDOWN_AVAILABLE:
            raise RuntimeError("Markdown is not installed. Please install markdown package.")
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        with open(md_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        html_content = markdown.markdown(
            md_content,
            extensions=['tables', 'fenced_code', 'codehilite']
        )
        
        styled_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; 
                       line-height: 1.6; max-width: 800px; margin: 40px auto; padding: 20px; }}
                h1, h2, h3 {{ color: #333; }}
                code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
                pre {{ background: #f4f4f4; padding: 15px; border-radius: 5px; overflow-x: auto; }}
                table {{ border-collapse: collapse; width: 100%; }}
                th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                th {{ background: #f4f4f4; }}
            </style>
        </head>
        <body>{html_content}</body>
        </html>
        """
        
        HTML(string=styled_html).write_pdf(str(output_path))
        return output_path
    
    @staticmethod
    def word_to_pdf(docx_path: Path, output_path: Path) -> Path:
        """Convert Word document to PDF (basic conversion via HTML)"""
        if not PYTHON_DOCX_AVAILABLE:
            raise RuntimeError("python-docx is not installed. Please install python-docx package.")
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        doc = WordDocument(str(docx_path))
        
        paragraphs_html = []
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text = para.text
            
            if "Heading 1" in style:
                paragraphs_html.append(f"<h1>{text}</h1>")
            elif "Heading 2" in style:
                paragraphs_html.append(f"<h2>{text}</h2>")
            elif "Heading 3" in style:
                paragraphs_html.append(f"<h3>{text}</h3>")
            else:
                paragraphs_html.append(f"<p>{text}</p>")
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; 
                       line-height: 1.6; margin: 40px; }}
                h1 {{ font-size: 24px; margin-bottom: 10px; }}
                h2 {{ font-size: 20px; margin-bottom: 8px; }}
                p {{ margin-bottom: 12px; }}
            </style>
        </head>
        <body>{''.join(paragraphs_html)}</body>
        </html>
        """
        
        HTML(string=html_content).write_pdf(str(output_path))
        return output_path
    
    @staticmethod
    def excel_to_pdf(xlsx_path: Path, output_path: Path) -> Path:
        """Convert Excel to PDF (basic conversion via HTML table)"""
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl is not installed. Please install openpyxl package.")
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        wb = openpyxl.load_workbook(str(xlsx_path))
        sheets_html = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows_html = []
            
            for row in sheet.iter_rows():
                cells = [f"<td>{cell.value or ''}</td>" for cell in row]
                rows_html.append(f"<tr>{''.join(cells)}</tr>")
            
            sheets_html.append(f"<h2>{sheet_name}</h2><table>{''.join(rows_html)}</table>")
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 20px; }}
                h2 {{ color: #333; }}
                table {{ border-collapse: collapse; width: 100%; margin-bottom: 30px; }}
                td, th {{ border: 1px solid #ddd; padding: 8px; font-size: 12px; }}
            </style>
        </head>
        <body>{''.join(sheets_html)}</body>
        </html>
        """
        
        HTML(string=html_content).write_pdf(str(output_path))
        return output_path
    
    @staticmethod
    def ppt_to_pdf(pptx_path: Path, output_path: Path) -> Path:
        """Convert PowerPoint to PDF (placeholder - requires LibreOffice for full conversion)"""
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed. Please install python-pptx package.")
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Please install weasyprint package.")
        
        prs = Presentation(str(pptx_path))
        slides_html = []
        
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(f"<p>{shape.text}</p>")
            
            slides_html.append(f"""
                <div class="slide">
                    <h2>Slide {i}</h2>
                    {''.join(texts)}
                </div>
            """)
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; }}
                .slide {{ page-break-after: always; padding: 40px; background: #f5f5f5; 
                         margin: 20px; border-radius: 10px; }}
                h2 {{ color: #333; border-bottom: 2px solid #007bff; padding-bottom: 10px; }}
            </style>
        </head>
        <body>{''.join(slides_html)}</body>
        </html>
        """
        
        HTML(string=html_content).write_pdf(str(output_path))
        return output_path
