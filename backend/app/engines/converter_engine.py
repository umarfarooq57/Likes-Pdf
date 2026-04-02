"""
Converter Engine - Complete Implementation
Handles all document format conversions:
- PDF ↔ Word, Excel, PowerPoint
- PDF ↔ Images
- PDF ↔ HTML, Markdown
- PDF → Text, CSV, XML, JSON
- PDF/A conversion
"""

import io
import os
import re
import json
import csv
import zipfile
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple, Union
from datetime import datetime

# Import processing libraries
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    PYMUPDF_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    Image = None
    PIL_AVAILABLE = False

try:
    from docx import Document as WordDocument
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    WordDocument = None
    DOCX_AVAILABLE = False

try:
    from openpyxl import Workbook, load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    Workbook = None
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    xlrd = None
    XLRD_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    markdown = None
    MARKDOWN_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Preformatted, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


class ConverterEngine:
    """
    Complete document conversion engine.
    Handles all format conversions.
    """

    # ==================== PDF to Images ====================

    @staticmethod
    def pdf_to_images(
        pdf_path: Path,
        output_dir: Path,
        format: str = "png",
        dpi: int = 150,
        pages: Optional[List[int]] = None,
        single_file: bool = False
    ) -> List[Path]:
        """
        Convert PDF pages to images.

        Args:
            pdf_path: Input PDF path
            output_dir: Output directory
            format: Image format (png, jpg, webp)
            dpi: Resolution (72-600)
            pages: Specific pages (1-indexed) or None for all
            single_file: If True, create single tall image

        Returns:
            List of image file paths
        """
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to image conversion")

        output_dir.mkdir(parents=True, exist_ok=True)
        doc = fitz.open(str(pdf_path))
        image_files = []
        stem = pdf_path.stem

        # DPI to zoom factor
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)

        # Determine pages to convert
        pages_to_convert = pages or list(range(1, len(doc) + 1))

        if single_file and PIL_AVAILABLE:
            # Create single concatenated image
            images = []
            for page_num in pages_to_convert:
                if 0 < page_num <= len(doc):
                    page = doc[page_num - 1]
                    pix = page.get_pixmap(matrix=matrix)
                    img = Image.frombytes(
                        "RGB", [pix.width, pix.height], pix.samples)
                    images.append(img)

            if images:
                # Calculate total height
                total_height = sum(img.height for img in images)
                max_width = max(img.width for img in images)

                # Create combined image
                combined = Image.new("RGB", (max_width, total_height), "white")
                y_offset = 0
                for img in images:
                    combined.paste(img, (0, y_offset))
                    y_offset += img.height

                output_path = output_dir / f"{stem}_combined.{format}"
                combined.save(str(output_path), format.upper())
                image_files.append(output_path)
        else:
            # Create individual images
            for page_num in pages_to_convert:
                if 0 < page_num <= len(doc):
                    page = doc[page_num - 1]
                    pix = page.get_pixmap(matrix=matrix)

                    output_path = output_dir / \
                        f"{stem}_page_{page_num}.{format}"
                    pix.save(str(output_path))
                    image_files.append(output_path)

        doc.close()
        return image_files

    @staticmethod
    def pdf_to_jpg(pdf_path: Path, output_dir: Path, quality: int = 85) -> List[Path]:
        """Convert PDF to JPG with quality setting"""
        return ConverterEngine.pdf_to_images(pdf_path, output_dir, format="jpg", dpi=150)

    @staticmethod
    def pdf_to_png(pdf_path: Path, output_dir: Path) -> List[Path]:
        """Convert PDF to PNG (lossless)"""
        return ConverterEngine.pdf_to_images(pdf_path, output_dir, format="png", dpi=150)

    # ==================== Images to PDF ====================

    @staticmethod
    def images_to_pdf(
        image_paths: List[Path],
        output_path: Path,
        page_size: str = "A4",
        margin: float = 0.5,
        fit_mode: str = "contain"
    ) -> Path:
        """
        Convert images to PDF.

        Args:
            image_paths: List of image file paths
            output_path: Output PDF path
            page_size: A4, Letter, Legal, or custom (width, height)
            margin: Margin in inches
            fit_mode: contain (fit inside), cover (fill page), stretch

        Returns:
            Output PDF path
        """
        if not PYMUPDF_AVAILABLE or not PIL_AVAILABLE:
            raise RuntimeError(
                "PyMuPDF and PIL required for image to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Page sizes in points (72 points = 1 inch)
        page_sizes = {
            "A4": (595, 842),
            "A3": (842, 1191),
            "Letter": (612, 792),
            "Legal": (612, 1008),
            "A5": (420, 595),
        }

        target_width, target_height = page_sizes.get(
            page_size, page_sizes["A4"])
        margin_pts = margin * 72

        doc = fitz.open()

        for img_path in image_paths:
            try:
                img = Image.open(str(img_path))

                # Convert to RGB if needed
                if img.mode in ("RGBA", "P", "LA"):
                    background = Image.new("RGB", img.size, "white")
                    if img.mode == "P":
                        img = img.convert("RGBA")
                    background.paste(img, mask=img.split()
                                     [-1] if img.mode == "RGBA" else None)
                    img = background
                elif img.mode != "RGB":
                    img = img.convert("RGB")

                img_width, img_height = img.size

                # Calculate available space
                available_width = target_width - 2 * margin_pts
                available_height = target_height - 2 * margin_pts

                # Calculate scaling based on fit mode
                if fit_mode == "contain":
                    scale = min(available_width / img_width,
                                available_height / img_height)
                elif fit_mode == "cover":
                    scale = max(available_width / img_width,
                                available_height / img_height)
                else:  # stretch
                    scale_x = available_width / img_width
                    scale_y = available_height / img_height
                    scale = min(scale_x, scale_y)

                new_width = img_width * scale
                new_height = img_height * scale

                # Create page
                page = doc.new_page(width=target_width, height=target_height)

                # Center image
                x = (target_width - new_width) / 2
                y = (target_height - new_height) / 2

                rect = fitz.Rect(x, y, x + new_width, y + new_height)

                # Save image to bytes
                img_bytes = io.BytesIO()
                img.save(img_bytes, format="PNG")
                img_bytes.seek(0)

                page.insert_image(rect, stream=img_bytes.getvalue())

            except Exception as e:
                print(f"Error processing image {img_path}: {e}")
                continue

        doc.save(str(output_path))
        doc.close()
        return output_path

    # ==================== PDF to Word ====================

    @staticmethod
    def pdf_to_word(
        pdf_path: Path,
        output_path: Path,
        preserve_layout: bool = True
    ) -> Path:
        """
        Convert PDF to Word document.
        Extracts text and attempts to preserve formatting.
        """
        if not PYMUPDF_AVAILABLE or not DOCX_AVAILABLE:
            raise RuntimeError(
                "PyMuPDF and python-docx required for PDF to Word conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))
        word_doc = WordDocument()

        for page_num, page in enumerate(doc):
            if page_num > 0:
                word_doc.add_page_break()

            if preserve_layout:
                # Extract text blocks with positions
                blocks = page.get_text("dict")["blocks"]

                for block in blocks:
                    if block["type"] == 0:  # Text block
                        for line in block.get("lines", []):
                            text = ""
                            for span in line.get("spans", []):
                                text += span.get("text", "")

                            if text.strip():
                                para = word_doc.add_paragraph(text)

                                # Try to apply formatting
                                if line.get("spans"):
                                    span = line["spans"][0]
                                    font_size = span.get("size", 12)

                                    for run in para.runs:
                                        run.font.size = Pt(font_size)

                    elif block["type"] == 1:  # Image block
                        # Extract and embed image
                        try:
                            xref = block.get("xref")
                            if xref:
                                img = doc.extract_image(xref)
                                if img:
                                    img_bytes = io.BytesIO(img["image"])
                                    word_doc.add_picture(
                                        img_bytes, width=Inches(5))
                        except Exception:
                            pass
            else:
                # Simple text extraction
                text = page.get_text()
                word_doc.add_paragraph(text)

        doc.close()
        word_doc.save(str(output_path))
        return output_path

    # ==================== Word to PDF ====================

    @staticmethod
    def word_to_pdf(
        docx_path: Path,
        output_path: Path
    ) -> Path:
        """
        Convert Word document to PDF.
        Uses python-docx and reportlab for conversion.
        """
        if not DOCX_AVAILABLE or not REPORTLAB_AVAILABLE:
            raise RuntimeError(
                "python-docx and reportlab required for Word to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Read Word document
        word_doc = WordDocument(str(docx_path))

        # Create PDF
        pdf_doc = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = getSampleStyleSheet()
        story = []

        for para in word_doc.paragraphs:
            text = para.text
            if text.strip():
                # Determine style based on paragraph style
                style_name = para.style.name if para.style else "Normal"

                if "Heading 1" in style_name:
                    p_style = styles["Heading1"]
                elif "Heading 2" in style_name:
                    p_style = styles["Heading2"]
                elif "Heading 3" in style_name:
                    p_style = styles["Heading3"]
                else:
                    p_style = styles["Normal"]

                story.append(Paragraph(text, p_style))
                story.append(Spacer(1, 12))

        if story:
            pdf_doc.build(story)
        else:
            # Create empty PDF
            c = canvas.Canvas(str(output_path), pagesize=A4)
            c.showPage()
            c.save()

        return output_path

    # ==================== PDF to Excel ====================

    @staticmethod
    def pdf_to_excel(
        pdf_path: Path,
        output_path: Path,
        extract_tables: bool = True
    ) -> Path:
        """
        Convert PDF tables to Excel.
        Extracts tabular data from PDF.
        """
        if not PYMUPDF_AVAILABLE or not OPENPYXL_AVAILABLE:
            raise RuntimeError(
                "PyMuPDF and openpyxl required for PDF to Excel conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Data"

        current_row = 1

        for page_num, page in enumerate(doc):
            # Add page header
            ws.cell(row=current_row, column=1,
                    value=f"--- Page {page_num + 1} ---")
            current_row += 1

            if extract_tables:
                # Try to extract tables using text analysis
                tables = page.find_tables()

                if tables:
                    for table in tables:
                        for row in table.extract():
                            for col_idx, cell in enumerate(row, 1):
                                ws.cell(row=current_row, column=col_idx,
                                        value=str(cell) if cell else "")
                            current_row += 1
                        current_row += 1  # Space between tables
                else:
                    # Fallback: extract text line by line
                    text = page.get_text()
                    for line in text.split("\n"):
                        if line.strip():
                            # Try to split by common delimiters
                            cells = re.split(r'\s{2,}|\t', line)
                            for col_idx, cell in enumerate(cells, 1):
                                ws.cell(row=current_row, column=col_idx,
                                        value=cell.strip())
                            current_row += 1
            else:
                # Simple text extraction
                text = page.get_text()
                for line in text.split("\n"):
                    if line.strip():
                        ws.cell(row=current_row, column=1, value=line)
                        current_row += 1

            current_row += 1  # Space between pages

        doc.close()
        wb.save(str(output_path))
        return output_path

    # ==================== Excel to PDF ====================

    @staticmethod
    def excel_to_pdf(
        xlsx_path: Path,
        output_path: Path
    ) -> Path:
        """Convert Excel to PDF"""
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError(
                "reportlab required for Excel to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        c = canvas.Canvas(str(output_path), pagesize=A4)
        width, height = A4

        file_ext = xlsx_path.suffix.lower()

        def draw_rows(sheet_name: str, rows):
            y = height - 50
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, y, f"Sheet: {sheet_name}")
            y -= 30

            c.setFont("Helvetica", 10)

            for row in rows:
                if y < 50:
                    c.showPage()
                    y = height - 50
                    c.setFont("Helvetica", 10)

                x = 50
                for cell in row:
                    value = "" if cell is None else str(cell)
                    c.drawString(x, y, value[:30])
                    x += 80
                    if x > width - 50:
                        break

                y -= 15

            c.showPage()

        # .xlsx/.xlsm path via openpyxl
        if file_ext in {".xlsx", ".xlsm"}:
            # True .xlsx files are zip containers; if not, try legacy fallback.
            if OPENPYXL_AVAILABLE and zipfile.is_zipfile(str(xlsx_path)):
                try:
                    wb = load_workbook(
                        str(xlsx_path), data_only=True, read_only=True)
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        draw_rows(sheet, ws.iter_rows(values_only=True))
                    c.save()
                    return output_path
                except Exception:
                    pass

            # Fallback: some files are old .xls content with .xlsx extension.
            xlrd_error = None
            if XLRD_AVAILABLE:
                try:
                    book = xlrd.open_workbook(str(xlsx_path))
                    for sheet in book.sheets():
                        rows = []
                        for r in range(sheet.nrows):
                            rows.append(sheet.row_values(r))
                        draw_rows(sheet.name, rows)
                    c.save()
                    return output_path
                except Exception as e:
                    xlrd_error = str(e)

            # Last-resort fallback: treat file as delimited text (e.g., CSV renamed as .xlsx).
            try:
                raw = xlsx_path.read_bytes()
                decoded = None
                for enc in ("utf-8", "utf-8-sig", "latin-1"):
                    try:
                        decoded = raw.decode(enc)
                        break
                    except Exception:
                        continue

                if decoded is not None:
                    lines = [ln for ln in decoded.splitlines() if ln.strip()]
                    if lines:
                        candidates = [",", "\t", ";", "|"]
                        delimiter = max(candidates, key=lambda d: sum(
                            ln.count(d) for ln in lines[:20]))
                        parsed_rows = [ln.split(delimiter) for ln in lines]
                        draw_rows("Recovered Data", parsed_rows)
                        c.save()
                        return output_path
            except Exception:
                pass

            if xlrd_error:
                raise RuntimeError(
                    "Failed to read Excel workbook. The file appears corrupted or unsupported. "
                    f"Details: {xlrd_error}"
                )

            raise RuntimeError(
                "openpyxl and xlrd are required for robust .xlsx conversion")

        # Legacy .xls path via xlrd
        if file_ext == ".xls":
            if not XLRD_AVAILABLE:
                raise RuntimeError("xlrd is required for .xls conversion")

            try:
                book = xlrd.open_workbook(str(xlsx_path))
                for sheet in book.sheets():
                    rows = []
                    for r in range(sheet.nrows):
                        rows.append(sheet.row_values(r))
                    draw_rows(sheet.name, rows)
            except Exception as e:
                raise RuntimeError(
                    "Failed to convert .xls file. Ensure xlrd is installed and workbook is valid. "
                    f"Details: {str(e)}"
                )

            c.save()
            return output_path

        raise RuntimeError(
            "Unsupported Excel format. Please use .xlsx, .xlsm, or .xls")

    # ==================== PDF to PowerPoint ====================

    @staticmethod
    def pdf_to_ppt(
        pdf_path: Path,
        output_path: Path
    ) -> Path:
        """
        Convert PDF to PowerPoint.
        Each page becomes a slide with the page as an image.
        """
        if not PYMUPDF_AVAILABLE or not PPTX_AVAILABLE:
            raise RuntimeError(
                "PyMuPDF and python-pptx required for PDF to PPT conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))
        prs = Presentation()

        # Set slide dimensions to match PDF aspect ratio
        if len(doc) > 0:
            page = doc[0]
            rect = page.rect
            aspect = rect.width / rect.height
            prs.slide_width = PptxInches(10)
            prs.slide_height = PptxInches(10 / aspect)

        for page_num, page in enumerate(doc):
            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)
                                  )  # 2x zoom for quality
            img_bytes = pix.tobytes("png")

            # Add slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add image
            img_stream = io.BytesIO(img_bytes)
            left = PptxInches(0)
            top = PptxInches(0)
            slide.shapes.add_picture(
                img_stream, left, top, width=prs.slide_width)

        doc.close()
        prs.save(str(output_path))
        return output_path

    # ==================== PowerPoint to PDF ====================

    @staticmethod
    def ppt_to_pdf(
        pptx_path: Path,
        output_path: Path
    ) -> Path:
        """Convert PowerPoint to PDF"""
        if not PPTX_AVAILABLE or not REPORTLAB_AVAILABLE:
            raise RuntimeError(
                "python-pptx and reportlab required for PPT to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation(str(pptx_path))

        c = canvas.Canvas(str(output_path), pagesize=A4)
        width, height = A4

        for slide_num, slide in enumerate(prs.slides, 1):
            y = height - 50
            c.setFont("Helvetica-Bold", 12)
            c.drawString(50, y, f"Slide {slide_num}")
            y -= 30

            c.setFont("Helvetica", 10)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text
                        if text.strip():
                            if y < 50:
                                c.showPage()
                                y = height - 50
                            c.drawString(50, y, text[:100])
                            y -= 15

            c.showPage()

        c.save()
        return output_path

    # ==================== PDF to Text ====================

    @staticmethod
    def pdf_to_text(
        pdf_path: Path,
        output_path: Path,
        preserve_layout: bool = False
    ) -> Path:
        """Extract text from PDF to TXT file"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to text conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))
        text_parts = []

        for page_num, page in enumerate(doc, 1):
            text_parts.append(f"\n{'='*50}\nPage {page_num}\n{'='*50}\n")

            if preserve_layout:
                text_parts.append(page.get_text("text"))
            else:
                text_parts.append(page.get_text())

        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(text_parts))

        return output_path

    # ==================== PDF to Text ====================
    @staticmethod
    def pdf_to_text(pdf_path: Path, output_path: Path) -> Path:
        """Convert PDF to plain text"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to text conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc = fitz.open(str(pdf_path))
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(text_parts))

        return output_path

    # ==================== PDF to HTML ====================

    @staticmethod
    def pdf_to_html(
        pdf_path: Path,
        output_path: Path,
        include_images: bool = True
    ) -> Path:
        """Convert PDF to HTML"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to HTML conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))

        html_parts = ["""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Converted PDF</title>
    <style>
        body {
            font-family: Arial, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
            line-height: 1.6;
        }
        .page {
            border-bottom: 2px solid #ccc;
            padding-bottom: 20px;
            margin-bottom: 20px;
        }
        .page-number {
            color: #666;
            font-size: 0.8em;
            margin-bottom: 10px;
        }
        img {
            max-width: 100%;
            height: auto;
        }
    </style>
</head>
<body>
"""]

        for page_num, page in enumerate(doc, 1):
            html_parts.append(f'<div class="page">')
            html_parts.append(
                f'<div class="page-number">Page {page_num}</div>')

            # Get HTML content
            html = page.get_text("html")
            html_parts.append(html)

            html_parts.append('</div>')

        html_parts.append("</body></html>")

        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_parts))

        return output_path

    # ==================== HTML to PDF ====================

    @staticmethod
    def html_to_pdf(
        html_content: Optional[str] = None,
        html_path: Optional[Path] = None,
        url: Optional[str] = None,
        output_path: Path = None,
        options: Dict[str, Any] = None
    ) -> Path:
        """Convert HTML to PDF"""
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("reportlab required for HTML to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Get HTML content
        if html_content:
            content = html_content
        elif html_path:
            with open(html_path, "r", encoding="utf-8") as f:
                content = f.read()
        elif url:
            import httpx
            response = httpx.get(url, timeout=30)
            content = response.text
        else:
            raise ValueError("Must provide html_content, html_path, or url")

        # Simple HTML to PDF using reportlab
        # For complex HTML, consider using weasyprint or playwright

        # Strip HTML tags for simple conversion
        from html.parser import HTMLParser

        class HTMLStripper(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []

            def handle_data(self, data):
                self.text.append(data)

            def get_text(self):
                return ''.join(self.text)

        stripper = HTMLStripper()
        stripper.feed(content)
        text = stripper.get_text()

        # Create PDF
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        for line in text.split('\n'):
            if line.strip():
                story.append(Paragraph(line, styles["Normal"]))
                story.append(Spacer(1, 6))

        if story:
            pdf_doc.build(story)
        else:
            c = canvas.Canvas(str(output_path), pagesize=A4)
            c.showPage()
            c.save()

        return output_path

    # ==================== Markdown to PDF ====================

    @staticmethod
    def markdown_to_pdf(
        markdown_content: str,
        output_path: Path,
        options: Dict[str, Any] = None
    ) -> Path:
        """Convert Markdown to PDF"""
        if not MARKDOWN_AVAILABLE or not REPORTLAB_AVAILABLE:
            raise RuntimeError(
                "markdown and reportlab required for Markdown to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Convert markdown to HTML
        html = markdown.markdown(
            markdown_content,
            extensions=['tables', 'fenced_code', 'codehilite']
        )

        # Convert HTML to PDF
        return ConverterEngine.html_to_pdf(html_content=html, output_path=output_path)

    # ==================== PDF to CSV ====================

    @staticmethod
    def pdf_to_csv(
        pdf_path: Path,
        output_path: Path
    ) -> Path:
        """Extract tabular data from PDF to CSV"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to CSV conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))
        rows = []

        for page in doc:
            tables = page.find_tables()

            if tables:
                for table in tables:
                    for row in table.extract():
                        rows.append(
                            [str(cell) if cell else "" for cell in row])

        doc.close()

        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(rows)

        return output_path

    # ==================== PDF to XML ====================

    @staticmethod
    def pdf_to_xml(
        pdf_path: Path,
        output_path: Path
    ) -> Path:
        """Convert PDF to XML format"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to XML conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))

        xml_parts = ['<?xml version="1.0" encoding="UTF-8"?>\n<document>']

        for page_num, page in enumerate(doc, 1):
            xml_parts.append(f'  <page number="{page_num}">')

            # Get text blocks
            blocks = page.get_text("dict")["blocks"]

            for block in blocks:
                if block["type"] == 0:  # Text
                    xml_parts.append('    <block type="text">')
                    for line in block.get("lines", []):
                        text = "".join(span.get("text", "")
                                       for span in line.get("spans", []))
                        # Escape XML special characters
                        text = text.replace("&", "&amp;").replace(
                            "<", "&lt;").replace(">", "&gt;")
                        xml_parts.append(f'      <line>{text}</line>')
                    xml_parts.append('    </block>')

            xml_parts.append('  </page>')

        xml_parts.append('</document>')

        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(xml_parts))

        return output_path

    # ==================== PDF to JSON ====================

    @staticmethod
    def pdf_to_json(
        pdf_path: Path,
        output_path: Path
    ) -> Path:
        """Convert PDF to JSON format"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to JSON conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))

        document_data = {
            "metadata": doc.metadata,
            "page_count": len(doc),
            "pages": []
        }

        for page_num, page in enumerate(doc, 1):
            page_data = {
                "number": page_num,
                "width": page.rect.width,
                "height": page.rect.height,
                "text": page.get_text(),
                "blocks": []
            }

            blocks = page.get_text("dict")["blocks"]

            for block in blocks:
                if block["type"] == 0:  # Text
                    block_data = {
                        "type": "text",
                        "bbox": block["bbox"],
                        "lines": []
                    }

                    for line in block.get("lines", []):
                        line_text = "".join(span.get("text", "")
                                            for span in line.get("spans", []))
                        block_data["lines"].append(line_text)

                    page_data["blocks"].append(block_data)

            document_data["pages"].append(page_data)

        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(document_data, f, indent=2, ensure_ascii=False)

        return output_path

    # ==================== PDF/A Conversion ====================

    @staticmethod
    def convert_to_pdfa(
        pdf_path: Path,
        output_path: Path,
        conformance: str = "PDF/A-2b"
    ) -> Path:
        """
        Convert PDF to PDF/A archival format.
        Note: Full PDF/A compliance requires specialized tools.
        This provides basic conversion.
        """
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF/A conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = fitz.open(str(pdf_path))

        # Embed fonts and linearize for PDF/A compatibility
        doc.save(
            str(output_path),
            garbage=4,
            deflate=True,
            clean=True,
            linear=True
        )

        doc.close()
        return output_path

    # ==================== Batch Conversion ====================

    @staticmethod
    def batch_convert(
        input_paths: List[Path],
        output_dir: Path,
        target_format: str,
        options: Dict[str, Any] = None
    ) -> List[Path]:
        """
        Batch convert multiple files.

        Args:
            input_paths: List of input file paths
            output_dir: Output directory
            target_format: Target format (pdf, png, jpg, docx, etc.)
            options: Format-specific options

        Returns:
            List of output file paths
        """
        output_dir.mkdir(parents=True, exist_ok=True)
        output_files = []
        options = options or {}

        for input_path in input_paths:
            try:
                stem = input_path.stem
                ext = input_path.suffix.lower()[1:]

                # Determine conversion method
                if ext == "pdf":
                    if target_format in ["png", "jpg", "jpeg", "webp"]:
                        files = ConverterEngine.pdf_to_images(
                            input_path, output_dir, format=target_format
                        )
                        output_files.extend(files)
                    elif target_format in ["docx", "doc"]:
                        output_path = output_dir / f"{stem}.docx"
                        ConverterEngine.pdf_to_word(input_path, output_path)
                        output_files.append(output_path)
                    elif target_format in ["xlsx", "xls"]:
                        output_path = output_dir / f"{stem}.xlsx"
                        ConverterEngine.pdf_to_excel(input_path, output_path)
                        output_files.append(output_path)
                    elif target_format in ["txt", "text"]:
                        output_path = output_dir / f"{stem}.txt"
                        ConverterEngine.pdf_to_text(input_path, output_path)
                        output_files.append(output_path)
                    elif target_format == "html":
                        output_path = output_dir / f"{stem}.html"
                        ConverterEngine.pdf_to_html(input_path, output_path)
                        output_files.append(output_path)

                elif ext in ["png", "jpg", "jpeg", "webp", "gif", "bmp"]:
                    if target_format == "pdf":
                        output_path = output_dir / f"{stem}.pdf"
                        ConverterEngine.images_to_pdf(
                            [input_path], output_path)
                        output_files.append(output_path)

                elif ext in ["docx", "doc"]:
                    if target_format == "pdf":
                        output_path = output_dir / f"{stem}.pdf"
                        ConverterEngine.word_to_pdf(input_path, output_path)
                        output_files.append(output_path)

                elif ext in ["xlsx", "xls"]:
                    if target_format == "pdf":
                        output_path = output_dir / f"{stem}.pdf"
                        ConverterEngine.excel_to_pdf(input_path, output_path)
                        output_files.append(output_path)

                elif ext in ["pptx", "ppt"]:
                    if target_format == "pdf":
                        output_path = output_dir / f"{stem}.pdf"
                        ConverterEngine.ppt_to_pdf(input_path, output_path)
                        output_files.append(output_path)

            except Exception as e:
                print(f"Error converting {input_path}: {e}")
                continue

        return output_files

    # ==================== PDF to CSV ====================

    @staticmethod
    def pdf_to_csv(pdf_path: Path, output_path: Path) -> Path:
        """Extract tables from PDF to CSV"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to CSV conversion")

        doc = fitz.open(str(pdf_path))
        all_data = []

        for page in doc:
            tabs = page.find_tables()
            for tab in tabs:
                all_data.extend(tab.extract())

        doc.close()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(all_data)

        return output_path

    # ==================== PDF to XML ====================

    @staticmethod
    def pdf_to_xml(pdf_path: Path, output_path: Path) -> Path:
        """Convert PDF structure to XML"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to XML conversion")

        doc = fitz.open(str(pdf_path))

        # Simple XML construction
        xml_content = ['<?xml version="1.0" encoding="UTF-8"?>']
        xml_content.append('<document name="{}">'.format(pdf_path.name))

        for i, page in enumerate(doc):
            xml_content.append('  <page number="{}">'.format(i+1))
            blocks = page.get_text("dict")["blocks"]
            for b in blocks:
                if b["type"] == 0:  # Text block
                    xml_content.append('    <block type="text">')
                    for l in b["lines"]:
                        for s in l["spans"]:
                            text = s["text"].replace("&", "&amp;").replace(
                                "<", "&lt;").replace(">", "&gt;")
                            xml_content.append(
                                '      <span>{}</span>'.format(text))
                    xml_content.append('    </block>')
            xml_content.append('  </page>')

        xml_content.append('</document>')
        doc.close()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(xml_content))

        return output_path

    # ==================== PDF to JSON ====================

    @staticmethod
    def pdf_to_json(pdf_path: Path, output_path: Path) -> Path:
        """Convert PDF structure to JSON"""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF required for PDF to JSON conversion")

        doc = fitz.open(str(pdf_path))
        data = {
            "filename": pdf_path.name,
            "page_count": len(doc),
            "pages": []
        }

        for i, page in enumerate(doc):
            page_data = {
                "number": i + 1,
                "width": page.rect.width,
                "height": page.rect.height,
                "content": page.get_text("dict")
            }
            data["pages"].append(page_data)

        doc.close()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2)

        return output_path

    # ==================== CSV to PDF ====================

    @staticmethod
    def csv_to_pdf(csv_path: Path, output_path: Path) -> Path:
        """Convert CSV to PDF report"""
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("reportlab required for CSV to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()

        data = []
        with open(csv_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                data.append(row)

        if data:
            # Adjust column widths based on data
            col_widths = None  # Auto
            t = Table(data, hAlign='CENTER')
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
            ]))
            elements.append(t)
        else:
            elements.append(Paragraph("Empty CSV file", styles['Normal']))

        doc.build(elements)
        return output_path

    # ==================== JSON to PDF ====================

    @staticmethod
    def json_to_pdf(json_path: Path, output_path: Path) -> Path:
        """Convert JSON to PDF report"""
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("reportlab required for JSON to PDF conversion")

        output_path.parent.mkdir(parents=True, exist_ok=True)

        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()

        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        elements.append(Paragraph("JSON Document Export", styles['Title']))
        elements.append(Paragraph(
            f"Generated at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
        elements.append(Spacer(1, 12))

        # Format JSON with indentation
        json_str = json.dumps(data, indent=4)

        # Use Preformatted for code-like layout
        elements.append(Preformatted(
            json_str, styles['BodyText'], maxLineLength=80))

        doc.build(elements)
        return output_path

    @staticmethod
    def zip_files(files: List[Path], output_path: Path) -> Path:
        """Zip multiple files into a single archive"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for file_path in files:
                if file_path.exists():
                    zipf.write(file_path, arcname=file_path.name)
        return output_path
